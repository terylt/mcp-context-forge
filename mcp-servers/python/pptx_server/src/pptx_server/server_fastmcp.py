#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Location: ./mcp-servers/python/pptx_server/src/pptx_server/server_fastmcp.py
Copyright 2025
SPDX-License-Identifier: Apache-2.0
Authors: Mihai Criveti

PowerPoint FastMCP Server

Comprehensive MCP server for creating and editing PowerPoint presentations.
Supports slide creation, text formatting, shapes, images, tables, and charts.
Powered by FastMCP for enhanced type safety and automatic validation.
"""

import logging
import sys
import uuid
from pathlib import Path
from typing import Any

from fastmcp import FastMCP
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pydantic import Field

# Configure logging to stderr to avoid MCP protocol interference
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    handlers=[logging.StreamHandler(sys.stderr)],
)
logger = logging.getLogger(__name__)

# Create FastMCP server instance
mcp = FastMCP("pptx-server")


class PresentationManager:
    """Manages PowerPoint presentations and operations."""

    def __init__(self):
        """Initialize the presentation manager."""
        self.presentations: dict[str, Presentation] = {}
        self.work_dir = Path("/tmp/pptx_server")
        self.work_dir.mkdir(exist_ok=True)

    def create_presentation(
        self, title: str | None = None, subtitle: str | None = None
    ) -> dict[str, Any]:
        """Create a new PowerPoint presentation."""
        try:
            prs = Presentation()
            pres_id = str(uuid.uuid4())

            # Add title slide if title provided
            if title or subtitle:
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                if title:
                    slide.shapes.title.text = title
                if subtitle and len(slide.placeholders) > 1:
                    slide.placeholders[1].text = subtitle

            self.presentations[pres_id] = prs

            # Save to file
            file_path = self.work_dir / f"{pres_id}.pptx"
            prs.save(str(file_path))

            return {
                "success": True,
                "presentation_id": pres_id,
                "file_path": str(file_path),
                "slide_count": len(prs.slides),
                "message": "Presentation created successfully",
            }
        except Exception as e:
            logger.error(f"Error creating presentation: {e}")
            return {"success": False, "error": str(e)}

    def add_slide(
        self,
        presentation_id: str,
        layout_index: int = 1,
        title: str | None = None,
        content: str | None = None,
    ) -> dict[str, Any]:
        """Add a new slide to the presentation."""
        try:
            if presentation_id not in self.presentations:
                return {"success": False, "error": "Presentation not found"}

            prs = self.presentations[presentation_id]

            # Ensure layout index is valid
            if layout_index >= len(prs.slide_layouts):
                layout_index = 1  # Default to content layout

            slide = prs.slides.add_slide(prs.slide_layouts[layout_index])

            # Set title if provided
            if title and slide.shapes.title:
                slide.shapes.title.text = title

            # Set content if provided
            if content:
                # Find content placeholder
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1:  # Content placeholder
                        shape.text = content
                        break

            # Save presentation
            file_path = self.work_dir / f"{presentation_id}.pptx"
            prs.save(str(file_path))

            return {
                "success": True,
                "slide_index": len(prs.slides) - 1,
                "total_slides": len(prs.slides),
                "message": "Slide added successfully",
            }
        except Exception as e:
            logger.error(f"Error adding slide: {e}")
            return {"success": False, "error": str(e)}

    def set_slide_title(self, presentation_id: str, slide_index: int, title: str) -> dict[str, Any]:
        """Set the title of a specific slide."""
        try:
            if presentation_id not in self.presentations:
                return {"success": False, "error": "Presentation not found"}

            prs = self.presentations[presentation_id]

            if slide_index >= len(prs.slides):
                return {"success": False, "error": "Slide index out of range"}

            slide = prs.slides[slide_index]

            if slide.shapes.title:
                slide.shapes.title.text = title
            else:
                return {"success": False, "error": "Slide has no title placeholder"}

            # Save presentation
            file_path = self.work_dir / f"{presentation_id}.pptx"
            prs.save(str(file_path))

            return {
                "success": True,
                "slide_index": slide_index,
                "title": title,
                "message": "Slide title updated successfully",
            }
        except Exception as e:
            logger.error(f"Error setting slide title: {e}")
            return {"success": False, "error": str(e)}

    def set_slide_content(
        self, presentation_id: str, slide_index: int, content: str
    ) -> dict[str, Any]:
        """Set the main content of a specific slide."""
        try:
            if presentation_id not in self.presentations:
                return {"success": False, "error": "Presentation not found"}

            prs = self.presentations[presentation_id]

            if slide_index >= len(prs.slides):
                return {"success": False, "error": "Slide index out of range"}

            slide = prs.slides[slide_index]

            # Find content placeholder
            content_set = False
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Content placeholder
                    shape.text = content
                    content_set = True
                    break

            if not content_set:
                return {"success": False, "error": "No content placeholder found"}

            # Save presentation
            file_path = self.work_dir / f"{presentation_id}.pptx"
            prs.save(str(file_path))

            return {
                "success": True,
                "slide_index": slide_index,
                "message": "Slide content updated successfully",
            }
        except Exception as e:
            logger.error(f"Error setting slide content: {e}")
            return {"success": False, "error": str(e)}

    def add_text_box(
        self,
        presentation_id: str,
        slide_index: int,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> dict[str, Any]:
        """Add a text box to a slide."""
        try:
            if presentation_id not in self.presentations:
                return {"success": False, "error": "Presentation not found"}

            prs = self.presentations[presentation_id]

            if slide_index >= len(prs.slides):
                return {"success": False, "error": "Slide index out of range"}

            slide = prs.slides[slide_index]

            # Add text box
            text_box = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            text_frame = text_box.text_frame
            text_frame.text = text

            # Save presentation
            file_path = self.work_dir / f"{presentation_id}.pptx"
            prs.save(str(file_path))

            return {
                "success": True,
                "slide_index": slide_index,
                "message": "Text box added successfully",
            }
        except Exception as e:
            logger.error(f"Error adding text box: {e}")
            return {"success": False, "error": str(e)}

    def add_image(
        self,
        presentation_id: str,
        slide_index: int,
        image_path: str,
        left: float,
        top: float,
        width: float | None = None,
        height: float | None = None,
    ) -> dict[str, Any]:
        """Add an image to a slide."""
        try:
            if presentation_id not in self.presentations:
                return {"success": False, "error": "Presentation not found"}

            prs = self.presentations[presentation_id]

            if slide_index >= len(prs.slides):
                return {"success": False, "error": "Slide index out of range"}

            if not Path(image_path).exists():
                return {"success": False, "error": "Image file not found"}

            slide = prs.slides[slide_index]

            # Add image
            if width and height:
                pic = slide.shapes.add_picture(
                    image_path, Inches(left), Inches(top), Inches(width), Inches(height)
                )
            elif width:
                pic = slide.shapes.add_picture(
                    image_path, Inches(left), Inches(top), width=Inches(width)
                )
            elif height:
                pic = slide.shapes.add_picture(
                    image_path, Inches(left), Inches(top), height=Inches(height)
                )
            else:
                pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top))

            # Save presentation
            file_path = self.work_dir / f"{presentation_id}.pptx"
            prs.save(str(file_path))

            return {
                "success": True,
                "slide_index": slide_index,
                "message": "Image added successfully",
            }
        except Exception as e:
            logger.error(f"Error adding image: {e}")
            return {"success": False, "error": str(e)}

    def add_shape(
        self,
        presentation_id: str,
        slide_index: int,
        shape_type: str,
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> dict[str, Any]:
        """Add a shape to a slide."""
        try:
            if presentation_id not in self.presentations:
                return {"success": False, "error": "Presentation not found"}

            prs = self.presentations[presentation_id]

            if slide_index >= len(prs.slides):
                return {"success": False, "error": "Slide index out of range"}

            slide = prs.slides[slide_index]

            # Map shape types
            shape_map = {
                "rectangle": MSO_SHAPE.RECTANGLE,
                "oval": MSO_SHAPE.OVAL,
                "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
                "diamond": MSO_SHAPE.DIAMOND,
                "star": MSO_SHAPE.STAR_5_POINT,
                "arrow": MSO_SHAPE.RIGHT_ARROW,
                "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
            }

            if shape_type not in shape_map:
                return {"success": False, "error": f"Unsupported shape type: {shape_type}"}

            # Add shape
            shape = slide.shapes.add_shape(
                shape_map[shape_type], Inches(left), Inches(top), Inches(width), Inches(height)
            )

            # Save presentation
            file_path = self.work_dir / f"{presentation_id}.pptx"
            prs.save(str(file_path))

            return {
                "success": True,
                "slide_index": slide_index,
                "shape_type": shape_type,
                "message": "Shape added successfully",
            }
        except Exception as e:
            logger.error(f"Error adding shape: {e}")
            return {"success": False, "error": str(e)}

    def add_table(
        self,
        presentation_id: str,
        slide_index: int,
        rows: int,
        cols: int,
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> dict[str, Any]:
        """Add a table to a slide."""
        try:
            if presentation_id not in self.presentations:
                return {"success": False, "error": "Presentation not found"}

            prs = self.presentations[presentation_id]

            if slide_index >= len(prs.slides):
                return {"success": False, "error": "Slide index out of range"}

            slide = prs.slides[slide_index]

            # Add table
            table = slide.shapes.add_table(
                rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
            ).table

            # Save presentation
            file_path = self.work_dir / f"{presentation_id}.pptx"
            prs.save(str(file_path))

            return {
                "success": True,
                "slide_index": slide_index,
                "rows": rows,
                "cols": cols,
                "message": "Table added successfully",
            }
        except Exception as e:
            logger.error(f"Error adding table: {e}")
            return {"success": False, "error": str(e)}

    def save_presentation(
        self, presentation_id: str, output_path: str | None = None
    ) -> dict[str, Any]:
        """Save the presentation to a file."""
        try:
            if presentation_id not in self.presentations:
                return {"success": False, "error": "Presentation not found"}

            prs = self.presentations[presentation_id]

            if output_path:
                file_path = Path(output_path)
            else:
                file_path = self.work_dir / f"{presentation_id}.pptx"

            # Ensure directory exists
            file_path.parent.mkdir(parents=True, exist_ok=True)

            # Save presentation
            prs.save(str(file_path))

            return {
                "success": True,
                "file_path": str(file_path),
                "file_size": file_path.stat().st_size,
                "slide_count": len(prs.slides),
                "message": "Presentation saved successfully",
            }
        except Exception as e:
            logger.error(f"Error saving presentation: {e}")
            return {"success": False, "error": str(e)}

    def get_presentation_info(self, presentation_id: str) -> dict[str, Any]:
        """Get information about a presentation."""
        try:
            if presentation_id not in self.presentations:
                return {"success": False, "error": "Presentation not found"}

            prs = self.presentations[presentation_id]

            slides_info = []
            for i, slide in enumerate(prs.slides):
                slide_info = {
                    "index": i,
                    "has_title": slide.shapes.title is not None,
                    "shape_count": len(slide.shapes),
                    "layout_name": slide.slide_layout.name,
                }
                if slide.shapes.title:
                    slide_info["title"] = slide.shapes.title.text
                slides_info.append(slide_info)

            return {
                "success": True,
                "presentation_id": presentation_id,
                "slide_count": len(prs.slides),
                "slides": slides_info,
                "slide_width": prs.slide_width,
                "slide_height": prs.slide_height,
            }
        except Exception as e:
            logger.error(f"Error getting presentation info: {e}")
            return {"success": False, "error": str(e)}

    def delete_slide(self, presentation_id: str, slide_index: int) -> dict[str, Any]:
        """Delete a slide from the presentation."""
        try:
            if presentation_id not in self.presentations:
                return {"success": False, "error": "Presentation not found"}

            prs = self.presentations[presentation_id]

            if slide_index >= len(prs.slides):
                return {"success": False, "error": "Slide index out of range"}

            # Remove slide from XML
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[slide_index])

            # Save presentation
            file_path = self.work_dir / f"{presentation_id}.pptx"
            prs.save(str(file_path))

            # Reload presentation to ensure consistency
            self.presentations[presentation_id] = Presentation(str(file_path))

            return {
                "success": True,
                "deleted_index": slide_index,
                "remaining_slides": len(self.presentations[presentation_id].slides),
                "message": "Slide deleted successfully",
            }
        except Exception as e:
            logger.error(f"Error deleting slide: {e}")
            return {"success": False, "error": str(e)}

    def open_presentation(self, file_path: str) -> dict[str, Any]:
        """Open an existing PowerPoint presentation."""
        try:
            if not Path(file_path).exists():
                return {"success": False, "error": "File not found"}

            prs = Presentation(file_path)
            pres_id = str(uuid.uuid4())
            self.presentations[pres_id] = prs

            return {
                "success": True,
                "presentation_id": pres_id,
                "slide_count": len(prs.slides),
                "message": "Presentation opened successfully",
            }
        except Exception as e:
            logger.error(f"Error opening presentation: {e}")
            return {"success": False, "error": str(e)}


# Initialize presentation manager
manager = PresentationManager()


# Tool definitions using FastMCP decorators
@mcp.tool(description="Create a new PowerPoint presentation")
async def create_presentation(
    title: str | None = Field(None, description="Title for the first slide"),
    subtitle: str | None = Field(None, description="Subtitle for the first slide"),
) -> dict[str, Any]:
    """Create a new PowerPoint presentation."""
    return manager.create_presentation(title, subtitle)


@mcp.tool(description="Open an existing PowerPoint presentation")
async def open_presentation(
    file_path: str = Field(..., description="Path to the PPTX file"),
) -> dict[str, Any]:
    """Open an existing PowerPoint presentation."""
    return manager.open_presentation(file_path)


@mcp.tool(description="Add a new slide to the presentation")
async def add_slide(
    presentation_id: str = Field(..., description="ID of the presentation"),
    layout_index: int = Field(1, ge=0, le=10, description="Slide layout index"),
    title: str | None = Field(None, description="Slide title"),
    content: str | None = Field(None, description="Slide content"),
) -> dict[str, Any]:
    """Add a new slide to the presentation."""
    return manager.add_slide(presentation_id, layout_index, title, content)


@mcp.tool(description="Set the title of a slide")
async def set_slide_title(
    presentation_id: str = Field(..., description="ID of the presentation"),
    slide_index: int = Field(..., ge=0, description="Index of the slide"),
    title: str = Field(..., description="New title for the slide"),
) -> dict[str, Any]:
    """Set the title of a specific slide."""
    return manager.set_slide_title(presentation_id, slide_index, title)


@mcp.tool(description="Set the main content of a slide")
async def set_slide_content(
    presentation_id: str = Field(..., description="ID of the presentation"),
    slide_index: int = Field(..., ge=0, description="Index of the slide"),
    content: str = Field(..., description="Content text for the slide"),
) -> dict[str, Any]:
    """Set the main content of a specific slide."""
    return manager.set_slide_content(presentation_id, slide_index, content)


@mcp.tool(description="Add a text box to a slide")
async def add_text_box(
    presentation_id: str = Field(..., description="ID of the presentation"),
    slide_index: int = Field(..., ge=0, description="Index of the slide"),
    text: str = Field(..., description="Text content"),
    left: float = Field(..., ge=0, le=10, description="Left position in inches"),
    top: float = Field(..., ge=0, le=10, description="Top position in inches"),
    width: float = Field(..., ge=0.1, le=10, description="Width in inches"),
    height: float = Field(..., ge=0.1, le=10, description="Height in inches"),
) -> dict[str, Any]:
    """Add a text box to a slide."""
    return manager.add_text_box(presentation_id, slide_index, text, left, top, width, height)


@mcp.tool(description="Add an image to a slide")
async def add_image(
    presentation_id: str = Field(..., description="ID of the presentation"),
    slide_index: int = Field(..., ge=0, description="Index of the slide"),
    image_path: str = Field(..., description="Path to the image file"),
    left: float = Field(..., ge=0, le=10, description="Left position in inches"),
    top: float = Field(..., ge=0, le=10, description="Top position in inches"),
    width: float | None = Field(None, ge=0.1, le=10, description="Width in inches"),
    height: float | None = Field(None, ge=0.1, le=10, description="Height in inches"),
) -> dict[str, Any]:
    """Add an image to a slide."""
    return manager.add_image(presentation_id, slide_index, image_path, left, top, width, height)


@mcp.tool(description="Add a shape to a slide")
async def add_shape(
    presentation_id: str = Field(..., description="ID of the presentation"),
    slide_index: int = Field(..., ge=0, description="Index of the slide"),
    shape_type: str = Field(
        ...,
        pattern="^(rectangle|oval|triangle|diamond|star|arrow|rounded_rectangle)$",
        description="Type of shape",
    ),
    left: float = Field(..., ge=0, le=10, description="Left position in inches"),
    top: float = Field(..., ge=0, le=10, description="Top position in inches"),
    width: float = Field(..., ge=0.1, le=10, description="Width in inches"),
    height: float = Field(..., ge=0.1, le=10, description="Height in inches"),
) -> dict[str, Any]:
    """Add a shape to a slide."""
    return manager.add_shape(presentation_id, slide_index, shape_type, left, top, width, height)


@mcp.tool(description="Add a table to a slide")
async def add_table(
    presentation_id: str = Field(..., description="ID of the presentation"),
    slide_index: int = Field(..., ge=0, description="Index of the slide"),
    rows: int = Field(..., ge=1, le=50, description="Number of rows"),
    cols: int = Field(..., ge=1, le=20, description="Number of columns"),
    left: float = Field(..., ge=0, le=10, description="Left position in inches"),
    top: float = Field(..., ge=0, le=10, description="Top position in inches"),
    width: float = Field(..., ge=0.1, le=10, description="Width in inches"),
    height: float = Field(..., ge=0.1, le=10, description="Height in inches"),
) -> dict[str, Any]:
    """Add a table to a slide."""
    return manager.add_table(presentation_id, slide_index, rows, cols, left, top, width, height)


@mcp.tool(description="Delete a slide from the presentation")
async def delete_slide(
    presentation_id: str = Field(..., description="ID of the presentation"),
    slide_index: int = Field(..., ge=0, description="Index of the slide to delete"),
) -> dict[str, Any]:
    """Delete a slide from the presentation."""
    return manager.delete_slide(presentation_id, slide_index)


@mcp.tool(description="Save the presentation to a file")
async def save_presentation(
    presentation_id: str = Field(..., description="ID of the presentation"),
    output_path: str | None = Field(None, description="Output file path"),
) -> dict[str, Any]:
    """Save the presentation to a file."""
    return manager.save_presentation(presentation_id, output_path)


@mcp.tool(description="Get information about the presentation")
async def get_presentation_info(
    presentation_id: str = Field(..., description="ID of the presentation"),
) -> dict[str, Any]:
    """Get information about a presentation."""
    return manager.get_presentation_info(presentation_id)


def main():
    """Main entry point for the FastMCP server."""
    import argparse

    parser = argparse.ArgumentParser(description="PowerPoint FastMCP Server")
    parser.add_argument(
        "--transport",
        choices=["stdio", "http"],
        default="stdio",
        help="Transport mode (stdio or http)",
    )
    parser.add_argument("--host", default="0.0.0.0", help="HTTP host")
    parser.add_argument("--port", type=int, default=9014, help="HTTP port")

    args = parser.parse_args()

    if args.transport == "http":
        logger.info(f"Starting PowerPoint FastMCP Server on HTTP at {args.host}:{args.port}")
        mcp.run(transport="http", host=args.host, port=args.port)
    else:
        logger.info("Starting PowerPoint FastMCP Server on stdio")
        mcp.run()


if __name__ == "__main__":
    main()
